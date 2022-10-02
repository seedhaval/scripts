from python_scripts.commonutil.helper import File
from pygments import highlight
from pygments.formatters import ImageFormatter
from pygments import lexers
from pptx import Presentation
from pptx.util import Inches


def main():
    height = 31
    max_scroll = 20
    buffer = 5
    width = 81

    fl = File(input("Enter full path to Python script : "))
    data = fl.readlines()
    data.extend([' '] * 10)
    data_fmt = [x + ' ' * (width - len(x)) + '.' for x in data]
    cnt = len(data)
    img_fldr = r"D:\tmp"
    start = 0
    ppt_fl = r"C:\Users\Dell\OneDrive\Desktop\code_walkthrough.pptx"
    tmplt_fl = r"C:\Users\Dell\OneDrive\Documents\template\code_walkthrough" \
               r".pptx"

    lexer = lexers.get_lexer_by_name("python", stripnl=False)

    prs = Presentation(tmplt_fl)
    title_slide_layout = prs.slide_layouts[0]

    for i, ln in enumerate(data, 1):
        ln_strp = ln.strip()
        ln = ln + (' ' * (width - len(ln)))
        if ln_strp.startswith('import ') or ln_strp.startswith('from '):
            pass
        elif ln_strp == '':
            pass
        else:
            if i - start >= max_scroll:
                start = i - 1 - buffer
                if start + height > cnt:
                    start = cnt - height
            code = '\n'.join(data_fmt[start:start + height])
            img_fl = img_fldr + "\\" + ("%03d" % i) + ".png"
            with open(img_fl, "wb") as f:
                f.write(highlight(code, lexer,
                                  ImageFormatter(line_number_chars=3,
                                                 line_number_start=start + 1,
                                                 font_size=28,
                                                 hl_lines=[i - start])))
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.add_picture(img_fl, left=0, top=0, width=Inches(13.4),
                                     height=Inches(7.5))

    prs.save(ppt_fl)

if __name__ == "__main__":
    main()